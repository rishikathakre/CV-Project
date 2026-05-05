from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

ROOT = Path(__file__).resolve().parent
path = ROOT / "Group19_Final_Presentation_v2.pptx"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
NEW_TEXT = (
    "System Weights:   α (dwell) = 0.30   |   β (revisits) = 0.30   |"
    "   γ (irregularity) = 0.20   |   δ (billing bypass) = 0.20"
    "   →   Grid Search Optimal F1 = 0.950"
)

prs  = Presentation(str(path))
sl5  = prs.slides[4]

OLD_FRAGMENT = "Optimal Weights"
patched = False
for shape in sl5.shapes:
    if not shape.has_text_frame:
        continue
    for para in shape.text_frame.paragraphs:
        full = "".join(r.text for r in para.runs)
        if OLD_FRAGMENT in full:
            for run in para.runs:
                run.text = ""
            if para.runs:
                para.runs[0].text = NEW_TEXT
                para.runs[0].font.size = Pt(12)
                para.runs[0].font.bold = True
                para.runs[0].font.color.rgb = NAVY
            patched = True
            print(f"Patched text box: {full[:60]}...")

if not patched:
    print("WARNING: target text box not found - nothing changed.")

prs.save(str(path))
print(f"Saved -> {path}")
